import os
import time
import shutil
import gc
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import qrcode
from PIL import Image
import comtypes.client
from multiprocessing import Pool


pd.options.mode.chained_assignment = None  # default='warn'


def load_config(file_path):
    with open(file_path, "r") as config_file:
        config = json.load(config_file)
    return config


base_path = os.path.dirname(os.path.abspath(__file__))
config = load_config("config.json")

# Extract configuration options
input_pptx = config["input_pptx"]
output_folder = config["output_folder"]
logo_path = config["logo_path"]
csv_file = config["csv_file"]


output_folder = os.path.join(base_path, output_folder)
pptx_output_folder = os.path.join(base_path, "certificates_pptx")
qr_code_folder = os.path.join(base_path, "qr_codes")


def PPT_to_PDF(input_pptx, output_pdf):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")

    if output_pdf[-3:] != "pdf":
        output_pdf = output_pdf + ".pdf"
    deck = powerpoint.Presentations.Open(input_pptx, WithWindow=False)
    deck.ExportAsFixedFormat(output_pdf, FixedFormatType=2)
    deck.Close()
    powerpoint.Quit()


def generate_qr_code(url, output_path):
    logo = Image.open(logo_path)
    basewidth = 100
    wpercent = basewidth / float(logo.size[0])
    hsize = int((float(logo.size[1]) * float(wpercent)))
    logo = logo.resize((basewidth, hsize), Image.Resampling.LANCZOS)
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=9,
        border=4,
    )
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color=(85, 101, 111), back_color="white").convert("RGB")
    pos = ((img.size[0] - logo.size[0]) // 2, (img.size[1] - logo.size[1]) // 2)
    img.paste(logo, pos)
    img = img.resize((100, 100), Image.Resampling.LANCZOS)
    img.save(output_path)


def process_pptx(row):
    os.makedirs(output_folder, exist_ok=True)
    os.makedirs(pptx_output_folder, exist_ok=True)
    os.makedirs(qr_code_folder, exist_ok=True)

    prs = Presentation(input_pptx)
    qr_code_url = row[3]
    qr_code_path = os.path.join(qr_code_folder, f"{row[1]}.png")
    generate_qr_code(qr_code_url, qr_code_path)

    # Define the placeholders to replace
    placeholders = {
        "Placeholder_Name": row[2],
        "Placeholder_refno": row[1],
    }

    # Iterate through slides and shapes to find and replace the placeholders
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, value in placeholders.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, value)
        left = Inches(9.5)
        top = Inches(7.4)
        qr_code = Image.open(qr_code_path)
        pic = slide.shapes.add_picture(qr_code_path, left, top, width=None, height=None)
        pic.click_action.hyperlink.address = qr_code_url

    # Save the modified PowerPoint presentation
    updated_pptx = os.path.join(pptx_output_folder, f"{row[1]}.pptx")
    prs.save(updated_pptx)

    # Convert the updated PowerPoint presentation to PDF
    output_pdf = os.path.join(output_folder, f"{row[1]}.pdf")
    PPT_to_PDF(updated_pptx, output_pdf)

    # print(f"Certificate for {row['name']} has been saved as '{output_pdf}'.")
    gc.collect()


def generate_certificates_summary(df):
    """
    Generate a Markdown file with a summary of the certificates generated.
    Can be used to validate the certificates generated.
    """
    certificates_md = "certificates.md"
    # Create a DataFrame for the Markdown table
    df_md = df[["id", "name", "profile_url"]]
    df_md["profile_url"] = df_md["profile_url"].apply(
        lambda x: f"[Link]({x})" if pd.notna(x) else ""
    )

    # Write the DataFrame to the Markdown file
    with open(certificates_md, mode="w", newline="") as md_file:
        md_file.write("# Certificates\n")
        time_now = time.strftime("%d %B %Y, %H:%M:%S", time.localtime())
        md_file.write(f"Generated on {time_now}\n\n")
        df_md.to_markdown(md_file, index=False)


def main():
    start_time = time.time()
    # Load data from CSV file
    df = pd.read_csv(csv_file)
    print("Generating Certificates...\nPlease wait...")

    # Use multiprocessing Pool to parallelize certificate generation
    with Pool() as pool:
        pool.map(process_pptx, df.itertuples(index=False, name=None))

    shutil.rmtree(pptx_output_folder)
    shutil.rmtree(qr_code_folder)
    generate_certificates_summary(df)

    time_elapsed = time.time() - start_time
    print(
        "Time elapsed for generating {} certificates: {:.2f} seconds.".format(
            len(df), time_elapsed
        )
    )
    print("All certificates have been generated.")
    gc.collect()


if __name__ == "__main__":
    main()
